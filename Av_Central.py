import requests, json, base64, time, os, collections
from collections import Counter
import matplotlib.pyplot as plt
import time,csv
from datetime import date
import pandas as pd
from pptx import Presentation
from pptx.util import Inches



class Authendication():
    def encode_credentials(self,username, password):
        cred = f'{username}:{password}'.encode()
        convert_base = base64.b64encode(cred).decode()
        return f'Basic {convert_base}'

    def get_authenticate(self,url, username, password):
        headers = {'Authorization': self.encode_credentials(username, password), 'accept': 'application/json'}
        login_res = requests.post(url, headers=headers)
        if login_res.status_code == 200:
            return login_res.json()['access_token']
        else: print(f'Not Able to Connect | Status Code : {login_res.status_code} | Reason : {login_res.reason}')



class Fetching_data():
    def fetching_alarams(self, url, token, days):
        headers = {'Authorization': f'Bearer {token}'}
        page_num=1
        dummy=100
        csv_list,customer_list,eventname_list = [],[],[]
        while True:


            data = {"page": page_num, "size":100, "find": {"alarm.suppressed": ["false"]},"sort": {"alarm.timestamp_occured": "desc"},"range": {"alarm.timestamp_occured": {"gte": f"now-{days}d", "lte": "now", "timeZone": "-0500"}}}
            response = requests.post(url=url, data=json.dumps(data), headers=headers)
            if response.status_code==200:
                data=response.json()
                for i in range(len(data['results'])):
                    alaram_dic={}
                    alaram_dic["Customer Name"]=data['results'][i]['tenantId'].replace("cn://","").replace(".alienvault.cloud","")
                    customer_list.append(data['results'][i]['tenantId'].replace("cn://","").replace(".alienvault.cloud",""))
                    alaram_dic["Event Type"]=data["results"][i]["alarm"]['rule_intent']
                    alaram_dic["Event Name"]=data["results"][i]["alarm"]["rule_method"]
                    eventname_list.append(data["results"][i]["alarm"]["rule_method"])
                    for j in range(len(data["results"][i]["events"])):
                        try:
                            alaram_dic["Sub user Name"]=data["results"][i]["events"][j]["source_username"]
                            alaram_dic["Target user Name"] = data["results"][i]["events"][j]["destination_username"]
                        except:pass
                        try:
                            alaram_dic["Raw log"] = data["results"][i]["events"][j]["log"]
                        except:pass
                    csv_list.append(alaram_dic)
                page_num += 1


                if dummy > int(data['total']): break

            else: print(f'Not Able to Connect | Status Code : {response.status_code} | Reason : {response.reason}')
            dummy += 100
            time.sleep(0.5)
        fields = ["Customer Name", "Event Type", "Event Name", "source Ip", "destination Ip", "Sub user Name",
                  "Target user Name", "Raw log"]
        customer_dic=dict(Counter(customer_list))
        eventname_dic=dict(Counter(eventname_list))
        ev_count=sorted(eventname_dic.items(), key=lambda kv: kv[1])[-10:]
        cu_count=sorted(customer_dic.items(), key=lambda kv: kv[1])[-10:]
        graph_data=[dict(ev_count),dict(cu_count)]

        return csv_list,fields,graph_data
    def featching_vulnerability(self,url,token,days):
        headers = {'Authorization': f'Bearer {token}'}
        page_num = 1
        dummy = 100
        csv_list = []
        while True:

            data = {"page": page_num, "size": 100, "find": {"vulnerability.isValid": ["false"]},
                    "sort": {"vulnerability.lastTimestamp": "desc"},
                    "range": {"vulnerability.lastTimestamp": {"gte": f"now-{days}d", "lte": "now", "timeZone": "-0500"}}}
            response = requests.post(url=url, data=json.dumps(data), headers=headers)
            if response.status_code == 200:
                data = response.json()
                for i in range(len(data['results'])):
                    vulnerability_dic = {}
                    vulnerability_dic["Customer Name"] = data['results'][i]['tenantId'].replace("cn://", "").replace(
                        ".alienvault.cloud", "")
                    vulnerability_dic["asset name"] = data["results"][i]["asset"]['name']
                    vulnerability_dic["vulnerability_name"] = data["results"][i]["vulnerability"]["name"]
                    vulnerability_dic["vulnerability_cve"] = data["results"][i]["vulnerability"]["cve"]
                    vulnerability_dic["vulnerability_score"] = data["results"][i]["vulnerability"]["cvssScore"]
                    vulnerability_dic["vulnerability_severity"] = data["results"][i]["vulnerability"]["cvssSeverity"]
                    vulnerability_dic["vulnerability_description"] = data["results"][i]["vulnerability"]["description"]
                    csv_list.append(vulnerability_dic)
                page_num += 1
                if dummy > int(data['total']): break

            else:
                print(f'Not Able to Connect | Status Code : {response.status_code} | Reason : {response.reason}')
            dummy += 100
            time.sleep(0.5)
        fields=["Customer Name","asset name","vulnerability_name","vulnerability_cve","vulnerability_score","vulnerability_severity","vulnerability_description"]

        return csv_list,fields



class Result_output():
    def csvfile(self,filepath,response,fields):

        with open(filepath,"w",newline="",encoding="utf-8")as new_csv:

            writer = csv.DictWriter(new_csv, fieldnames=fields)
            writer.writeheader()
            writer.writerows(response)
    def excelfile(self,path_1,path_2):
        df=pd.read_csv(path_1)
        df.to_excel(path_2,index=False)
class Graph():
    def bar_graph(self,x,y,title,xname,yname,path):
        plt.bar(x,y,align = "center",alpha=0.5,color = ["r","b","g","y"],linewidth=0,width=0.5)
        plt.ylabel(yname)
        plt.xlabel(xname)
        plt.xticks(rotation=90)
        today = date.today()
        plt.gca().set_title(title + " " + today.strftime("%d/%m/%Y"), pad=20)
        plt.gca().spines['top'].set_visible(False)
        plt.gca().spines['right'].set_visible(False)
        plt.gca().set_ylim([0, max(y)])
        for i, j in enumerate(y):
            plt.text(x=i - 0.1, y=j + 1, s=str(j))
        plt.savefig(path, bbox_inches="tight")
        plt.clf()
class PPT():
    def ppt(self,graph_data,title,fields,path):
        list_1, list_2 = [], []
        for k,v in graph_data[0].items():
            list_1.append(k)
            list_2.append(v)
        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes
        rows =5
        cols = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table

        shapes.title.text = title
        for i in range(1):
            for j in range(2):
                table.cell(i,j).text=fields[j]

        for i in range(1,len(list_1)):
            for j in range(2):
                table.cell(i,0).text=list_1[i-1]
                table.cell(i,1).text=list_2[i-1]
        prs.save(path)

outh = Authendication()
f_data=Fetching_data()
rop=Result_output()

bg=Graph()
pt=PPT()
token=outh.get_authenticate()
response,fields,graph_data=f_data.fetching_alarams( token=token,days="1")
rop.csvfile(response=response,fields=fields)
response,fields=f_data.featching_vulnerability(token=token)
rop.csvfile(response=response,fields=fields)
rop.excelfile()





